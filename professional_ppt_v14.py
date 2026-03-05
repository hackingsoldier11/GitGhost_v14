import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_professional_presentation():
    target_path = r"d:\PROJRCT  CAS\GitGhost\GitGhost_Final_v14\GitGhost_v14\PROJECT_V14_PRO.pptx"
    prs = Presentation()
    
    # Professional Dark Theme Colors (Matching GitGhost UI)
    bg_color = RGBColor(10, 11, 15)       # Very dark navy/black
    text_color = RGBColor(230, 237, 243)  # Off-white
    ghost_green = RGBColor(0, 255, 65)    # Matrix Green
    accent_blue = RGBColor(88, 166, 255)  # Soft Blue
    alert_red = RGBColor(255, 0, 85)      # Critical Red

    def format_title_shape(shape, color=ghost_green):
        shape.text_frame.paragraphs[0].font.color.rgb = color
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.size = Pt(44)

    def format_body_shape(shape):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.color.rgb = text_color
            paragraph.font.size = Pt(18)

    # --- Slide 1: Welcome / Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, bg_color)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "GitGhost v14.F\nFull Stack Forensic Platform"
    format_title_shape(title)
    
    subtitle.text = "Audit the Past. Protect the Future.\nProject Architect: Antigravity AI | Version 14.F Release"
    subtitle.text_frame.paragraphs[0].font.color.rgb = accent_blue

    # --- Slide 2: The Core Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "The 'Ghost' Vulnerability"
    format_title_shape(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Deletion is Not Destruction"
    p = tf.add_paragraph(); p.text = "Developer Fallacy: Deleting a file (git rm) removes the risk."; p.level = 1
    p = tf.add_paragraph(); p.text = "The Reality: Secrets persist in the Git DAG as unreferenced blobs."; p.level = 1
    p = tf.add_paragraph(); p.text = "The Attack Vector: Hackers walk history to exhume historical API keys."; p.level = 1
    format_body_shape(slide.placeholders[1])

    # --- Slide 3: Solution Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Hybrid Forensic Architecture"
    format_title_shape(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Dual-Vector Exhumation Engine"
    p = tf.add_paragraph(); p.text = "Vector 1 (Live Shield): Active HEAD workspace scanning."; p.level = 1
    p = tf.add_paragraph(); p.text = "Vector 2 (Ghost Exhumation): Topological history traversal."; p.level = 1
    p = tf.add_paragraph(); p.text = "Result: 100% coverage of both current and historical code leakage."; p.level = 1
    format_body_shape(slide.placeholders[1])

    # --- Slide 4: AI Remediation Advisor ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "AI Remediation Advisor"
    format_title_shape(slide.shapes.title, accent_blue)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Contextual Security Guidance"
    p = tf.add_paragraph(); p.text = "LLM Core: Understands findings based on CVSS and Entropy."; p.level = 1
    p = tf.add_paragraph(); p.text = "Direct Action: Generates specific 'git filter-repo' commands."; p.level = 1
    p = tf.add_paragraph(); p.text = "Interactive Triage: Integrated chatbot for forensic Q&A."; p.level = 1
    format_body_shape(slide.placeholders[1])

    # --- Slide 5: ML Anomaly Detection ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Isolation Forest Outlier Detection"
    format_title_shape(slide.shapes.title, alert_red)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Probabilistic Secret Identification"
    p = tf.add_paragraph(); p.text = "Unsupervised ML: Identifies anomalies without rigid signatures."; p.level = 1
    p = tf.add_paragraph(); p.text = "Feature Vector: Entropy Score + Sequence Length + CVSS Density."; p.level = 1
    p = tf.add_paragraph(); p.text = "Graphical Output: Live scatterplots highlighting red outlier nodes."; p.level = 1
    format_body_shape(slide.placeholders[1])

    # --- Slide 6: OWASP 2025 Compliance ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "OWASP 2025 Matrix & Intelligence"
    format_title_shape(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Forward-Looking Compliance"
    p = tf.add_paragraph(); p.text = "Web 2025: Mapping artifacts to W01 (Access) and W02 (Crypto)."; p.level = 1
    p = tf.add_paragraph(); p.text = "LLM 2025: Prevention of training data leakage (L02/L05 vectors)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Scorecarding: Real-time compliance health bars based on findings."; p.level = 1
    format_body_shape(slide.placeholders[1])

    # --- Slide 7: Integrated Dashboard ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Full Stack Command Center"
    format_title_shape(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "Live Forensic Feed: Visualizing the DAG exhumation status."; p.level = 1
    p = tf.add_paragraph(); p.text = "Learning Studio: Forensic modules for educational onboarding."; p.level = 1
    p = tf.add_paragraph(); p.text = "Timeline Analysis: Historical secret trends across 365+ days."; p.level = 1
    format_body_shape(slide.placeholders[1])

    # --- Slide 8: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, bg_color)
    title = slide.shapes.title
    title.text = "GitGhost v14.F\nProject Complete"
    format_title_shape(title)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Ready for Security Audit | SOC 2 Ready | Forensic Lab Compliant"
    subtitle.text_frame.paragraphs[0].font.color.rgb = ghost_green

    prs.save(target_path)
    print(f"Professional presentation saved to: {target_path}")

if __name__ == "__main__":
    create_professional_presentation()
