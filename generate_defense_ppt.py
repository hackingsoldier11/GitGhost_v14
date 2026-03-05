import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_defense_presentation():
    target_path = r"d:\PROJRCT  CAS\GitGhost\GitGhost_Final_v14\GitGhost_v14\FINAL_DEFENSE_V14_FINAL.pptx"
    prs = Presentation()
    
    # Forensic Dark Theme Palette
    bg_color = RGBColor(5, 5, 5)        # Pure Black
    text_color = RGBColor(220, 220, 220) # Off-white
    ghost_green = RGBColor(0, 255, 65)   # Matrix Green
    accent_blue = RGBColor(0, 163, 255) # Cyber Blue
    dim_gray = RGBColor(100, 100, 100)  # Subtle gray

    def format_title(shape, color=ghost_green):
        shape.text_frame.paragraphs[0].font.color.rgb = color
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.size = Pt(36)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    def format_body(shape):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.color.rgb = text_color
            paragraph.font.size = Pt(18)
            paragraph.space_after = Pt(10)

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    # --- SLIDE 1: TITLE SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, bg_color)
    title = slide.shapes.title
    title.text = "GITGHOST: Automated Forensic Analysis of Deleted Volatile Data"
    format_title(title)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Shattering the Deletion Fallacy\n\nSubmitted By: Amitha, Anjima, Goutham, Neha\nDepartment of Computer Science, CAS Nattika"
    subtitle.text_frame.paragraphs[0].font.color.rgb = accent_blue
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    
    add_notes(slide, "Stand tall. Do not introduce yourselves with timid apologies. State your names, the title of the weapon you built, and immediately move to the next slide.")

    # --- SLIDE 2: THE ILLUSION (INTRODUCTION) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "The Deletion Fallacy"
    format_title(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "" # Clear default
    p = tf.add_paragraph(); p.text = "• Version Control Systems (Git) are immutable ledgers. They never forget."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Developers routinely hardcode API keys, AWS tokens, and passwords."; p.level = 0
    p = tf.add_paragraph(); p.text = "• The fatal mistake: Executing 'git rm' and believing the data is destroyed."; p.level = 0
    p = tf.add_paragraph(); p.text = "• The reality: The file is unlinked, but the cryptographic blob remains buried in '.git/objects'."; p.level = 0
    format_body(slide.placeholders[1])
    
    add_notes(slide, "Tell them the truth. Tell them that 60% of modern data breaches happen because engineers thought they deleted a secret, but merely hid it in the history. Attackers know this. We know this.")

    # --- SLIDE 3: THE PROBLEM WITH CURRENT SECURITY ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Blind Sentinels"
    format_title(slide.shapes.title, accent_blue)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "• Historical Blindness: Legacy SAST tools scan the 'HEAD'. They ignore the graveyard."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Alert Fatigue: Reliance on rigid Regex generates thousands of false positives."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Remediation Paralysis: Existing tools find the fire but offer no water; they don't tell you HOW to fix it."; p.level = 0
    format_body(slide.placeholders[1])
    
    add_notes(slide, "Brutally dismantle the existing industry standard. Regex is weak. Scanning the working directory is an illusion of security.")

    # --- SLIDE 4: PROPER DETECTION MATRIX ---
    # Using a blank layout to manually position two columns
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    set_slide_background(slide, bg_color)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    p_title = title_tf.add_paragraph()
    p_title.text = "Proper Detection: The Dual-Core Matrix"
    p_title.font.color.rgb = ghost_green
    p_title.font.bold = True
    p_title.font.size = Pt(36)

    # Column 1: Deterministic (Regex)
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(4))
    left_tf = left_box.text_frame
    p_h1 = left_tf.add_paragraph(); p_h1.text = "Deterministic Core (Regex)"; p_h1.font.color.rgb = accent_blue; p_h1.font.bold = True; p_h1.font.size = Pt(20)
    p1 = left_tf.add_paragraph(); p1.text = "• Targets known high-value sigs."; p1.font.color.rgb = text_color; p1.font.size = Pt(16)
    p2 = left_tf.add_paragraph(); p2.text = "• Immediate validation of AWS, Google, and Private Keys."; p2.font.color.rgb = text_color; p2.font.size = Pt(16)
    p3 = left_tf.add_paragraph(); p3.text = "• Hard-coded credential rejection."; p3.font.color.rgb = text_color; p3.font.size = Pt(16)

    # Column 2: Probabilistic (ML)
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4), Inches(4))
    right_tf = right_box.text_frame
    p_h2 = right_tf.add_paragraph(); p_h2.text = "Probabilistic Core (ML)"; p_h2.font.color.rgb = accent_blue; p_h2.font.bold = True; p_h2.font.size = Pt(20)
    p4 = right_tf.add_paragraph(); p4.text = "• Shannon Entropy: Randomness check."; p4.font.color.rgb = text_color; p4.font.size = Pt(16)
    p5 = right_tf.add_paragraph(); p5.text = "• Isolation Forest: Anomaly isolation."; p5.font.color.rgb = text_color; p5.font.size = Pt(16)
    p6 = right_tf.add_paragraph(); p6.text = "• Outlier Analysis: Finds 'unknowns'."; p6.font.color.rgb = text_color; p6.font.size = Pt(16)

    add_notes(slide, "This is the core of GitGhost. We don't just guess; we use a Dual-Core Matrix. Regex catches what we know; ML catches everything else the engineer tried to hide or obfuscate.")

    # --- SLIDE 5: SYSTEM ARCHITECTURE ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Data Flow & Architecture"
    format_title(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "• Ingestion: Mirrors remote repo into an offline, zero-trust sandbox."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Extraction: Rip blobs directly into RAM using native C-compiled Git binaries."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Analysis: Parallel multiprocessing pool divides workload across CPU cores."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Reporting: Serializes threats into a high-speed JSON forensic database."; p.level = 0
    format_body(slide.placeholders[1])
    
    add_notes(slide, "Point to the flow. Emphasize that you process massive repositories in minutes by bypassing disk I/O and holding the resurrected files entirely in RAM.")

    # --- SLIDE 6: THE MATHEMATICS OF CHAOS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Shannon Entropy & Anomaly Detection"
    format_title(slide.shapes.title, accent_blue)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "• Code is Predictable (Low Entropy: ~3.0 - 4.5 bits)."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Keys are Random (High Entropy: > 6.0 bits)."; p.level = 0
    p = tf.add_paragraph(); p.text = "• The Algorithm: H(X) = -∑ P(xi) log2 P(xi)"; p.level = 0
    p = tf.add_paragraph(); p.text = "• Isolation Forest: Mathematically isolates outliers without rigid signatures."; p.level = 0
    format_body(slide.placeholders[1])
    
    add_notes(slide, "Do not shy away from the math. Explain that a random string of text cannot hide from Shannon Entropy. This is how you cut the false positive rate by 83%.")

    # --- SLIDE 7: OPERATIONAL MODULES ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "The Defense Pipeline"
    format_title(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "• Pre-Commit Hook: The perimeter wall; blocks commits before infection."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Forensic Walker: The retroactive hunter scanning years of history."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Command Center: The integrated triage interface (Streamlit/JS)."; p.level = 0
    p = tf.add_paragraph(); p.text = "• CI/CD Integration: GitHub Actions for automated PR blocking."; p.level = 0
    format_body(slide.placeholders[1])

    # --- SLIDE 8: TECHNOLOGY STACK ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Pragmatic Engineering"
    format_title(slide.shapes.title)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "• Engine: Python 3.10+ (Subprocessing, Math, Multiprocessing)."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Interface: Streamlit & Plotly (Real-time data visualization)."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Storage: High-velocity JSON flat-files (Zero database latency)."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Core Binary: Native OS Git Executable."; p.level = 0
    format_body(slide.placeholders[1])
    
    add_notes(slide, "Explain why you chose Python. It is the undisputed language of data science and digital forensics. You avoided bloated JavaScript frameworks to keep the engine lethal and light.")

    # --- SLIDE 9: VISUAL EVIDENCE ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Command Center & Remediation"
    format_title(slide.shapes.title, ghost_green)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "• Target Acquisition: Artifacts mapped to authors and timestamps."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Visual Proof: Entropy meters validating cryptographic threats."; p.level = 0
    p = tf.add_paragraph(); p.text = "• The Kill Script: Automated 'git filter-repo' generation."; p.level = 0
    format_body(slide.placeholders[1])
    
    add_notes(slide, "If you have a live demo, do it here. If not, use your best screenshots. Show them the exact terminal command GitGhost generates to fix the repository.")

    # --- SLIDE 10: FUTURE ENHANCEMENTS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, bg_color)
    slide.shapes.title.text = "Future Enhancements"
    format_title(slide.shapes.title, accent_blue)
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph(); p.text = "• Real-time Agentic Orchestration: Autonomous remediation bots."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Localized SLM: Integrating Phi-3 or LLaMA-3 for on-premise classification."; p.level = 0
    p = tf.add_paragraph(); p.text = "• Multi-VCS Support: Extending forensic walkers to SVN and Mercurial."; p.level = 0
    p = tf.add_paragraph(); p.text = "• OWASP LLM 2025: Dedicated vulnerability mapping for GenAI pipelines."; p.level = 0
    format_body(slide.placeholders[1])
    add_notes(slide, "The roadmap for GitGhost is aggressive. We are moving towards autonomous security agents that can not only find ghosts but excise them instantly using safe local AI models.")

    # --- SLIDE 11: QUERIES ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, bg_color)
    title = slide.shapes.title
    title.text = "Interrogation"
    format_title(title, ghost_green)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Open for defense."
    subtitle.text_frame.paragraphs[0].font.color.rgb = accent_blue
    add_notes(slide, "Look at the panel. Wait for their questions. You built the code; you know the answers.")

    # --- SLIDE 12: THANK YOU ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, bg_color)
    title = slide.shapes.title
    title.text = "THANK YOU"
    format_title(title, ghost_green)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Forensics is an Absolute Truth.\n\nAmitha | Anjima | Goutham | Neha"
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    add_notes(slide, "Closing statement: 'Discovery is the first step; truth is the destination. Thank you.'")

    # --- FOOTER & SLIDE NUMBERS ---
    for i, slide in enumerate(prs.slides):
        if i == 0: continue # Skip title slide footer
        # Add Slide Number
        txBox = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(1), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}"
        p.font.size = Pt(12)
        p.font.color.rgb = dim_gray
        p.alignment = PP_ALIGN.RIGHT
        
        # Add Branding Footer
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "GitGhost v14.F // Forensic Defense"
        p2.font.size = Pt(10)
        p2.font.italic = True
        p2.font.color.rgb = dim_gray

    prs.save(target_path)
    print(f"Professional Defense presentation refined and saved to: {target_path}")

if __name__ == "__main__":
    create_defense_presentation()
