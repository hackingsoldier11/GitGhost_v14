import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_project_presentation():
    # Ensure directory exists
    target_dir = r"d:\PROJRCT  CAS\GitGhost\GitGhost_Final_v14\GitGhost_v14"
    if not os.path.exists(target_dir):
        os.makedirs(target_dir)
    
    prs = Presentation()
    
    # Slide Layouts
    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    # 1. Title Slide: GitGhost v14.F Full Stack
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "GitGhost v14.F\nFull Stack Forensic Platform"
    subtitle.text = "Exhuming Secrets from the Git DAG\nHybrid Scanning | AI Advisor | OWASP 2025"
    
    # 2. Hybrid Scanning Architecture
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Hybrid Scanning Engine"
    tf = slide.placeholders[1].text_frame
    tf.text = "Dual-Vector Forensics"
    p = tf.add_paragraph(); p.text = "Live Shield: Real-time scan of current workspace (HEAD)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Ghost Exhumation: Deep walk through Git history (all commits/branches)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Reconstruction: Exhumes 'deleted' blobs using topological DAG analysis."; p.level = 1

    # 3. AI-Powered Remediation
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "AI Remediation Advisor"
    tf = slide.placeholders[1].text_frame
    tf.text = "Context-Aware Triage"
    p = tf.add_paragraph(); p.text = "Interactive Chat: Advisor provides custom 'git filter-repo' commands."; p.level = 1
    p = tf.add_paragraph(); p.text = "Dynamic Analysis: Explains risks based on Shannon Entropy & CVSS."; p.level = 1
    p = tf.add_paragraph(); p.text = "Interconnectivity: Click findings in dashboard to trigger AI guidance."; p.level = 1

    # 4. Intelligence & Compliance
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Forensic Threat Intel"
    tf = slide.placeholders[1].text_frame
    tf.text = "OWASP 2025 & Isolation Forest"
    p = tf.add_paragraph(); p.text = "OWASP 2025: Mapping findings to Web (W01-W10) and LLM (L01-L10) vectors."; p.level = 1
    p = tf.add_paragraph(); p.text = "ML Anomalies: Outlier detection using Isolation Forest (Entropy vs CVSS)."; p.level = 1
    p = tf.add_paragraph(); p.text = "Data Sovereignty: Prevents 'GHOST' secrets from poisoning AI training sets."; p.level = 1

    # 5. Full Stack Dashboard
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Command Center v14.F"
    tf = slide.placeholders[1].text_frame
    tf.text = "Integrated Visual Analytics"
    p = tf.add_paragraph(); p.text = "Live Forensic Feed: Real-time logging of exhumation telemetry."; p.level = 1
    p = tf.add_paragraph(); p.text = "Learning Studio: Interactive modules for Forensic Analysts."; p.level = 1
    p = tf.add_paragraph(); p.text = "Scorecarding: Overall security health based on CVSS density."; p.level = 1

    # 6. Conclusion
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "GitGhost: The Future of\nSource Control Forensics"
    subtitle = slide.placeholders[1]
    subtitle.text = "Version 14.F // FULL STACK READY"

    save_path = os.path.join(target_dir, "PROJECT.pptx")
    prs.save(save_path)
    print(f"Project presentation updated and saved to: {save_path}")

if __name__ == "__main__":
    create_project_presentation()
