import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 0 - Title, 1 - Title and Content
    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    # Theme Colors
    bg_color = RGBColor(10, 10, 10)       # Dark background
    text_color = RGBColor(230, 237, 243)  # Light gray/blue text
    ghost_green = RGBColor(0, 255, 65)    # Ghost Green
    
    # Create Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "GitGhost v14.0\nForensic Academy"
    subtitle.text = "Proper Learning Setup & Secret Exhumation\nAntigravity AI Project Architect"
    
    # Module 1: The Forensic Mindset
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Module 1: The Forensic Mindset"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Deletion Fallacy"
    p = tf.add_paragraph()
    p.text = "Deleting a file containing secrets does NOT destroy the data. It persists as a 'Ghost' blob in the Git object database."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "The Ghost Vector: Attackers walk the Git Directed Acyclic Graph (DAG) for deleted objects (`--diff-filter=D`)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Solution: Kinetic Exhumation - reconstructs blobs from the object database without checking out the code."
    p.level = 1
    
    # Module 2: The Math of Secrets
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Module 2: The Math of Secrets"
    tf = slide.placeholders[1].text_frame
    tf.text = "Shannon Entropy as a Signature"
    p = tf.add_paragraph()
    p.text = "Instead of relying purely on rigid Regex, GitGhost uses mathematics to detect cryptographic chaos."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Formula: H = -∑ P(x) log₂ P(x)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Threshold Context: Standard text clusters ~3.5 to 4.5. Base64 encoded RSA keys and secure secrets exceed 6.0 bits of entropy."
    p.level = 1

    # Module 3: Iso-Forest ML 
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Module 3: Isolation Forest ML"
    tf = slide.placeholders[1].text_frame
    tf.text = "Unsupervised Anomaly Detection"
    p = tf.add_paragraph()
    p.text = "Used to drastically reduce False Positives."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Maps variables like: Entropy Score, Sequence Length, and Context Specificity."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Outliers are easily identified because random secrets isolate faster than structured standard code, plotting as distinct red nodes in the telemetry scatterplot."
    p.level = 1

    # Module 4: Purge Protocol
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Module 4: The 3-Step Purge"
    tf = slide.placeholders[1].text_frame
    tf.text = "Remediation Mechanics"
    p = tf.add_paragraph()
    p.text = "Step 1: ROTATE. Immediately invalidate the exposed credentials in the cloud provider."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Step 2: REWRITE. Execute `git filter-repo` to systematically destroy the object hashes across the entire DAG timeline."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Step 3: ENFORCE. Force-push to remotes and activate GitGhost's pre-commit hook (`hooks/pre-commit.py`) to block future breaches."
    p.level = 1

    # Lab Section
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Interactive Labs"
    tf = slide.placeholders[1].text_frame
    tf.text = "Apply Forensic Techniques in Sandboxes"
    p = tf.add_paragraph()
    p.text = "Lab 1: Vulnerable Juice Shop. Clone the OWASP target and launch `python gitghost.py juice-shop` to triage 38+ secrets via the dashboard."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Lab 2: Pre-Commit Defense. Apply the GitGhost hook locally and attempt to commit an AWS Access Key. Observe the immediate hard rejection."
    p.level = 1

    # Conclusion
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "To protect the future,\nwe must audit the past."
    slide.placeholders[1].text = "GitGhost Forensic Team\nv14.F"

    prs.save("GitGhost_Learning_Setup.pptx")
    print("Successfully generated GitGhost_Learning_Setup.pptx")

if __name__ == "__main__":
    create_presentation()
